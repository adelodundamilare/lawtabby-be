
from decouple import config
import os
import fitz
import PyPDF2
from PIL import Image
from io import BytesIO
from reportlab.lib import colors
import pytesseract
import pdfkit
from docx import Document
import io

from PyPDF2 import PdfReader, PdfWriter
from django.core.files.base import ContentFile
from zipfile import ZipFile
from django.contrib.sites.shortcuts import get_current_site
from pptx import Presentation
import pandas as pd

from reportlab.pdfgen.canvas import Canvas
from reportlab.lib.pagesizes import letter, landscape, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Image, Spacer, Table, TableStyle
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import letter, landscape

from PIL import Image as PILImage

import openai

from .models import OcrPdf, ProtectedPDF,MergedPDF, CompressedPDF, SplitPDF, OrganizedPdf, StampPdf, UnlockPdf

from django.conf import settings

TEMP_PATH = settings.TEMP_PATH


def protect_pdf(request, input_file, password, user):
    pdf_reader = PdfReader(input_file)
    pdf_writer = PdfWriter()
    input_file_name = input_file.name


    for page_num in range(len(pdf_reader.pages)):
        pdf_writer.add_page(pdf_reader.pages[page_num])

    pdf_writer.encrypt(password)

    buffer = BytesIO()
    pdf_writer.write(buffer)

    protected_file = ProtectedPDF(user=user, password=password)
    protected_file.protected_file.save(input_file_name, ContentFile(buffer.getvalue()))
    protected_file.save()

    current_site = get_current_site(request)
    base_url = f'http://{current_site.domain}'
    full_url = f'{base_url}{protected_file.protected_file.url}'

    return protected_file, full_url


def merge_pdf(request, user, pdf_list):
    pdf_writer = PdfWriter()

    for pdf_file in pdf_list:
        pdf_reader = PdfReader(pdf_file)
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            pdf_writer.add_page(page)

    buffer = BytesIO()
    pdf_writer.write(buffer)

    merged_pdf = MergedPDF(user=user)
    merged_pdf.merged_file.save('merged_output.pdf', ContentFile(buffer.getvalue()))
    merged_pdf.save()

    current_site = get_current_site(request)
    base_url = f'http://{current_site.domain}'
    full_url = f'{base_url}{merged_pdf.merged_file.url}'

    return merged_pdf, full_url


def compress_pdf(request, user, input_pdf, compression_quality):
    try:
        # Save the uploaded PDF file to a temporary location

        # temp_file_path = os.path.join(temp_path, input_pdf.name)
        temp_file_path = os.path.join(TEMP_PATH, input_pdf.name)

        with open(temp_file_path, 'wb') as temp_file:
            for chunk in input_pdf.chunks():
                temp_file.write(chunk)

        # Create a PDF document object using PyMuPDF
        pdf_document = fitz.open(temp_file_path)
        pdf_writer = fitz.open()
        print(pdf_writer, 'pdf_writer')

        for page_number in range(pdf_document.page_count):
            page = pdf_document[page_number]


            # Convert the page to a Pixmap
            pixmap = page.get_pixmap()

            # Convert the Pixmap to a PIL Image
            pil_image = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)

            # Compress the PIL Image
            quality = get_compression_quality(compression_quality)
            compressed_image_stream = BytesIO()
            pil_image.save(compressed_image_stream, format="JPEG", quality=quality)

            # Create a new PDF document with the compressed image
            compressed_pdf_page = pdf_writer.new_page(width=pixmap.width, height=pixmap.height)
            compressed_pdf_page.insert_image(compressed_pdf_page.rect, stream=compressed_image_stream.getvalue())

        # Save the compressed PDF
        buffer = BytesIO()
        # pdf_writer.write(buffer)
        pdf_writer.save(buffer)

        # Create a new CompressedPDF model instance
        compressed_pdf = CompressedPDF(user=user, compression_quality=compression_quality)
        compressed_pdf.compressed_file.save('compressed_output.pdf', ContentFile(buffer.getvalue()))
        compressed_pdf.save()

        # Clean up the temporary file
        try:
            os.remove(temp_file_path)
        except OSError as e:
            print(f"Error deleting temporary file: {e}")

        current_site = get_current_site(request)
        base_url = f'http://{current_site.domain}'
        full_url = f'{base_url}{compressed_pdf.compressed_file.url}'



        return compressed_pdf, full_url

    except Exception as e:
        print(f"An error occurred: {str(e)}")
        return None

def stamp_pdf_with_text(input_pdf, stamp_text, user):
    try:
        temp_file_path = os.path.join(TEMP_PATH, input_pdf.name)

        with open(temp_file_path, 'wb') as temp_file:
            for chunk in input_pdf.chunks():
                temp_file.write(chunk)

        pdf_reader = PdfReader(temp_file_path)
        pdf_writer = PdfWriter()

        watermark_buffer = BytesIO()

        canvas = Canvas(watermark_buffer, pagesize=letter)
        canvas.setFillColor(colors.Color(0, 0, 0, alpha=0.2))
        canvas.setFont("Helvetica", 36)

        text_width = canvas.stringWidth(stamp_text, "Helvetica", 36)
        text_height = 36

        center_x = letter[0] / 2 - (text_width / 2)
        center_y = letter[1] / 2 - (text_height / 2)

        # Rotate the canvas and draw the rotated text
        # canvas.rotate(30)  # Rotate by 30 degrees clockwise
        canvas.drawString(center_x, center_y, stamp_text)
        canvas.save()

        for page_number in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_number]
            watermark_pdf = PdfReader(BytesIO(watermark_buffer.getvalue()))
            watermark_page = watermark_pdf.pages[0]
            page.merge_page(watermark_page)
            pdf_writer.add_page(page)

        output_buffer = BytesIO()
        pdf_writer.write(output_buffer)

        stamped_pdf_instance = StampPdf(user=user)
        stamped_pdf_instance.pdf.save('stamped_output.pdf', ContentFile(output_buffer.getvalue()))
        stamped_pdf_instance.save()

        return stamped_pdf_instance

    except Exception as e:
        print(f"An error occurred: {str(e)}")
        return None


def get_compression_quality(choice):
    # Define compression quality based on user choice
    if choice == 'extreme':
        return 100   # Less quality, high compression
    elif choice == 'recommended':
        return 50   # Good quality, good compression
    elif choice == 'less':
        return 10   # High quality, less compression
    else:
        raise ValueError("Invalid compression choice")




def split_pdf(request, input_pdf, start_page, end_page, user):
    pdf_reader = PdfReader(input_pdf)
    print(f'Total number of pages: {len(pdf_reader.pages)}')
    ranges = [(start_page, end_page)]

    for i, (start, end) in enumerate(ranges):
        pdf_writer = PdfWriter()

        for page_num in range(start, end + 1):
            pdf_writer.add_page(pdf_reader.pages[page_num])

        buffer = BytesIO()
        pdf_writer.write(buffer)

        # Create a new SplitPDF instance and save it
        split_pdf_instance = SplitPDF(start_page=start, end_page=end, user=user)
        split_pdf_instance.split_pdf.save('split.pdf', ContentFile(buffer.getvalue()))
        split_pdf_instance.save()

    current_site = get_current_site(request)
    base_url = f'http://{current_site.domain}'
    full_url = f'{base_url}{split_pdf_instance.split_pdf.url}'  # Adjust as needed

    return split_pdf_instance, full_url


#convert pdf to images

def convert_pdf_to_image(input_pdf):
    temp_file_path = os.path.join(TEMP_PATH, input_pdf.name)
    with open(temp_file_path, 'wb') as temp_file:
        for chunk in input_pdf.chunks():
            temp_file.write(chunk)

    with fitz.open(temp_file_path) as pdf_document:  # Use a context manager to ensure proper closing
        image_paths = []
        for page_number in range(pdf_document.page_count):
            page = pdf_document[page_number]
            image = page.get_pixmap()
            image_data = image.tobytes()  # Extract raw image bytes
            image_paths.append(image_data)  # Append raw bytes directly

    os.remove(temp_file_path)  # Clean up the temporary file
    return image_paths


def create_zip_file(images, user):
    zip_buffer = BytesIO()
    with ZipFile(zip_buffer, 'w') as zip_file:
        for i, image_data in enumerate(images):
            zip_file.writestr(f'page_{i + 1}.jpeg', image_data)

    # Ensure the directory exists before saving the zip file
    zip_dir = os.path.join('pdf_images/zips/', str(user.id))
    os.makedirs(zip_dir, exist_ok=True)

    zip_name = f'pdf_images.zip'  # Simplified zip file name
    zip_file_path = os.path.join(zip_dir, zip_name)
    with open(zip_file_path, 'wb') as zip_file:
        zip_file.write(zip_buffer.getvalue())

    return zip_file_path, zip_buffer.getvalue()



def save_uploaded_file(input_file):
    """
    Saves the uploaded file to a temporary location.
    """
    temp_file_path = os.path.join(TEMP_PATH, input_file.name)
    with open(temp_file_path, 'wb') as temp_file:
        for chunk in input_file.chunks():
            temp_file.write(chunk)
    return temp_file_path

def convert_docx_to_pdf(docx_file):
    try:
        doc = Document(docx_file)
        buffer = io.BytesIO()
        pdf = SimpleDocTemplate(buffer, pagesize=letter)
        elements = []
        styles = getSampleStyleSheet()

        for para in doc.paragraphs:
            elements.append(Paragraph(para.text, styles['Normal']))
            elements.append(Spacer(1, 12))

        pdf.build(elements)
        return buffer.getvalue()
    except Exception as e:
        print(f"Error converting Word document: {str(e)}")
        raise ValueError("Error converting Word document")

def convert_pptx_to_pdf(pptx_file):
    try:
        prs = Presentation(pptx_file)
        buffer = io.BytesIO()
        pdf = SimpleDocTemplate(buffer, pagesize=landscape(letter))
        elements = []
        styles = getSampleStyleSheet()

        for slide in prs.slides:
            elements.append(Paragraph(f"Slide {slide.slide_id}", styles['Heading1']))
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    elements.append(Paragraph(shape.text, styles['Normal']))
            elements.append(Spacer(1, 12))

        pdf.build(elements)
        return buffer.getvalue()
    except Exception as e:
        print(f"Error converting PowerPoint document: {str(e)}")
        raise ValueError("Error converting PowerPoint document")

def convert_excel_to_pdf(excel_file, file_extension):
    try:
        if file_extension == 'xlsx':
            df = pd.read_excel(excel_file, engine='openpyxl')
        elif file_extension == 'xls':
            df = pd.read_excel(excel_file, engine='xlrd')
        else:
            raise ValueError(f"Unsupported Excel format: {file_extension}")

        buffer = io.BytesIO()
        pdf = SimpleDocTemplate(buffer, pagesize=landscape(letter))
        elements = []
        styles = getSampleStyleSheet()

        elements.append(Paragraph(f"Sheet: {excel_file.name}", styles['Heading1']))
        data = [df.columns.tolist()] + df.values.tolist()
        table = Table(data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 14),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 0), (-1, -1), 12),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        elements.append(table)

        pdf.build(elements)
        return buffer.getvalue()
    except Exception as e:
        print(f"Error converting Excel document: {str(e)}")
        raise ValueError("Error converting Excel document")

def convert_csv_to_pdf(csv_file):
    try:
        df = pd.read_csv(csv_file)

        buffer = io.BytesIO()
        pdf = SimpleDocTemplate(buffer, pagesize=landscape(letter))
        elements = []
        styles = getSampleStyleSheet()

        elements.append(Paragraph(f"CSV: {csv_file.name}", styles['Heading1']))
        data = [df.columns.tolist()] + df.values.tolist()
        table = Table(data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 14),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 0), (-1, -1), 12),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        elements.append(table)

        pdf.build(elements)
        return buffer.getvalue()
    except Exception as e:
        print(f"Error converting CSV document: {str(e)}")
        raise ValueError("Error converting CSV document")

def convert_image_to_pdf(image_file):
    try:
        image = PILImage.open(image_file)

        buffer = io.BytesIO()
        pdf = SimpleDocTemplate(buffer, pagesize=A4)
        width, height = A4

        # Adjust image size to fit the page
        img_width, img_height = image.size
        aspect = img_height / float(img_width)

        if img_width > width:
            img_width = width
            img_height = aspect * img_width

        elements = [Image(image_file, width=img_width, height=img_height)]

        pdf.build(elements)
        return buffer.getvalue()
    except Exception as e:
        print(f"Error converting image: {str(e)}")
        raise ValueError("Error converting image")

def convert_other_to_pdf(file):
    file_extension = file.name.split('.')[-1].lower()

    if file_extension in ['docx', 'doc']:
        return convert_docx_to_pdf(file)
    elif file_extension in ['pptx', 'ppt']:
        return convert_pptx_to_pdf(file)
    elif file_extension in ['xlsx', 'xls']:
        return convert_excel_to_pdf(file)
    elif file_extension == 'csv':
        return convert_csv_to_pdf(file)
    elif file_extension in ['jpg', 'jpeg', 'png', 'gif']:
        return convert_image_to_pdf(file)
    else:
        raise ValueError(f"Unsupported file format: {file_extension}")



def organize_pdf(input_pdf, user_order, pages_to_exclude, user):

    if isinstance(user_order, str):
        user_order = list(map(int, user_order.split(',')))
    if isinstance(pages_to_exclude, str):
        pages_to_exclude = list(map(int, pages_to_exclude.split(',')))

    with input_pdf.open(mode='rb') as pdf_file:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        total_pages = len(pdf_reader.pages)

        remaining_pages = [page for page in range(1, total_pages + 1) if page not in pages_to_exclude]

        if sorted(user_order) != sorted(remaining_pages):
            raise ValueError("Invalid page order. Please enter a valid order.")

        pdf_writer = PyPDF2.PdfWriter()
        for page_number in user_order:
            pdf_writer.add_page(pdf_reader.pages[page_number - 1])

        buffer = BytesIO()
        pdf_writer.write(buffer)
        buffer.seek(0)

        organized_pdf_instance = OrganizedPdf(user=user)
        organized_pdf_instance.organize_pdf.save(f"organized_output.pdf", ContentFile(buffer.getvalue()))
        organized_pdf_instance.save()

        print("PDF successfully organized.")
        return organized_pdf_instance


def unlock_pdf(input_pdf, password, user):
    pdf_reader = PyPDF2.PdfReader(input_pdf)


    # Try to decrypt the PDF with the provided password
    success = pdf_reader.decrypt(password)

    if success:
        # Create a PDF writer object
        pdf_writer = PyPDF2.PdfWriter()

        # Add each page from the decrypted PDF to the writer
        for page_num in range(len(pdf_reader.pages)):
            pdf_writer.add_page(pdf_reader.pages[page_num])

        # Create a BytesIO buffer and write the PDF content
        buffer = BytesIO()
        pdf_writer.write(buffer)
        buffer.seek(0)

        # Save the unlocked PDF to the UnlockPdf model
        unlock_pdf_instance = UnlockPdf(user=user)
        unlock_pdf_instance.unlock_pdf.save(f"unlocked_output.pdf", ContentFile(buffer.getvalue()))
        unlock_pdf_instance.save()

        print("PDF unlocked successfully.")
        return unlock_pdf_instance
    else:
        raise ValueError("Failed to unlock the PDF. Incorrect password.")


def pdf_to_ocr(input_pdf, user):
    temp_file_path = os.path.join(TEMP_PATH, input_pdf.name)

    with open(temp_file_path, 'wb') as temp_file:
        for chunk in input_pdf.chunks():
            temp_file.write(chunk)

    if not os.path.exists(temp_file_path):
        raise FileNotFoundError(f"Input PDF file '{temp_file_path}' not found.")

    # Configure Tesseract OCR path
    pytesseract.pytesseract.tesseract_cmd = "C:\\Program Files\\Tesseract-OCR\\tesseract.exe"

    # Check Tesseract OCR version
    if not pytesseract.pytesseract.get_tesseract_version():
        raise RuntimeError("Tesseract OCR is not properly configured. Please set the correct path.")

    pdf_document = fitz.open(temp_file_path)
    pdf_writer = fitz.open()

    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        image_bytes = page.get_pixmap().tobytes()
        image = Image.open(BytesIO(image_bytes))
        ocr_text = pytesseract.image_to_string(image)

        # Add OCR text to a new page in the output PDF
        new_page = pdf_writer.new_page(width=page.rect.width, height=page.rect.height)
        new_page.insert_text((0, 0), ocr_text)

    # Save the OCR result as a PDF
    buffer = BytesIO()
    pdf_writer.save(buffer)
    buffer.seek(0)

    # Save the PDF to the database using the OcrPdf model (adjust this part according to your model)
    pdf = OcrPdf(user=user)
    pdf.pdf.save('ocr_output.pdf', ContentFile(buffer.getvalue()))
    pdf.save()

    return pdf

def convert_pdf_page_to_image(pdf_document, page_num):
    page = pdf_document.load_page(page_num)
    pixmap = page.get_pixmap()
    image = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
    return image

def perform_ocr_on_image(image):
    ocr_dict = pytesseract.image_to_data(image, lang='eng', output_type=pytesseract.Output.DICT)
    ocr_text = " ".join([word for word in ocr_dict['text'] if word.strip()])
    return ocr_text

def summarize_pdf(input_pdf, user):
    temp_file_path = open_file(input_pdf)
    document = read_pdf(temp_file_path)

    chunks = split_text(document)

    summaries = []
    for chunk in chunks:
        prompt = "Please summarize the following document: \n"
        prompt += "Kindly ensure response is properly formatted using appropriate html tags (heading, paragraph and list tags only): \n"
        summary = gpt3_completion(prompt + chunk)

        if summary.startswith("GPT-3 error"):
            continue

        summaries.append(summary)

    # pdf_writer = fitz.open()

    # # Save the result as a PDF
    # buffer = BytesIO()
    # pdf_writer.save(buffer)
    # buffer.seek(0)

    # # Save the PDF to the database using the OcrPdf model (adjust this part according to your model)
    # pdf = PdfModel(user=user)
    # pdf.pdf.save('output.pdf', ContentFile(buffer.getvalue()))
    # pdf.save()

    return ''.join(summaries)

def read_pdf(filename):
    context = ""

    with fitz.open(filename) as pdf_file:
        num_pages = pdf_file.page_count
        for page_num in range(num_pages):
            page = pdf_file[page_num]
            page_text = page.get_text()
            context += page_text

    return context

def open_file(input_pdf):
    temp_file_path = os.path.join(TEMP_PATH, input_pdf.name)

    with open(temp_file_path, 'wb') as temp_file:
        for chunk in input_pdf.chunks():
            temp_file.write(chunk)

    if not os.path.exists(temp_file_path):
        raise FileNotFoundError(f"Input PDF file '{temp_file_path}' not found.")

    return temp_file_path

def split_text(text, chunk_size=5000):
    """
    Splits the given text into chunks of approximately the specified chunk size.

    Args:
        text (str): The text to split.
        chunk_size (int): The desired size of each chunk (in characters).

    Returns:
        List[str]: A list of chunks, each of approximately the specified chunk size.
    """
    chunks = []
    current_chunk = ""
    current_size = 0

    # Split the text into words
    words = text.split()

    for word in words:
        word_size = len(word) + 1  # Add 1 for the space between words

        # If adding the current word exceeds the chunk size, start a new chunk
        if current_size + word_size > chunk_size:
            chunks.append(current_chunk.strip())  # Add the current chunk to the list of chunks
            current_chunk = ""  # Reset the current chunk
            current_size = 0  # Reset the current size

        # Add the current word to the current chunk
        current_chunk += word + " "
        current_size += word_size

    # Add the remaining text as the last chunk
    if current_chunk:
        chunks.append(current_chunk.strip())

    return chunks


def gpt3_completion(prompt):
    prompt = prompt.encode(encoding='ASCII',errors='ignore').decode()
    try:
        client = openai.OpenAI(
            api_key=config('OPENAI_API_KEY')
        )

        response = client.chat.completions.create(
            messages=[
                {
                    "role": "user",
                    "content": prompt,
                }
            ],
            model="gpt-3.5-turbo",
        )
        return response.choices[0].message.content
    except Exception as oops:
        return "GPT-3 error: %s" % oops